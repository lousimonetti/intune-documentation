from __future__ import annotations

import json
from dataclasses import asdict, replace
from pathlib import Path
from typing import Dict, Iterable

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import BubbleChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from .reports.schema import RenderedReport


SECTION_KEYS = {
    "summary": "summary",
    "assets": "assets",
    "assignment_coverage": "assignment_coverage",
}


def _filter_sections(report: RenderedReport, include_sections: Iterable[str]) -> RenderedReport:
    allowed = {SECTION_KEYS.get(section, section) for section in include_sections}
    if not allowed:
        return report

    filtered_sections = [
        section
        for section in report.sections
        if any(key in allowed for key in section.payload.keys())
    ]
    return replace(report, sections=filtered_sections)


def write_rendered_reports(
    rendered: Dict[str, RenderedReport],
    output_prefix: Path,
    include_sections: Iterable[str],
) -> Dict[str, Path]:
    output_paths: Dict[str, Path] = {}
    output_prefix.parent.mkdir(parents=True, exist_ok=True)

    for format_name, report in rendered.items():
        filtered_report = _filter_sections(report, include_sections)
        output_path = _write_report_output(filtered_report, output_prefix, format_name)
        output_paths[format_name] = output_path

    return output_paths


def _write_report_output(report: RenderedReport, output_prefix: Path, format_name: str) -> Path:
    json_output_path = output_prefix.with_name(f"{output_prefix.name}-{format_name}.json")
    json_output_path.write_text(
        json.dumps(asdict(report), indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    if format_name == "word":
        docx_output_path = output_prefix.with_name(f"{output_prefix.name}-{format_name}.docx")
        _write_docx_report(report, docx_output_path)
        return docx_output_path
    if format_name == "ppt":
        pptx_output_path = output_prefix.with_name(f"{output_prefix.name}-{format_name}.pptx")
        _write_pptx_report(report, pptx_output_path)
        return pptx_output_path
    if format_name == "excel":
        xlsx_output_path = output_prefix.with_name(f"{output_prefix.name}-{format_name}.xlsx")
        _write_excel_report(report, xlsx_output_path)
        return xlsx_output_path
    return json_output_path


def _write_docx_report(report: RenderedReport, output_path: Path) -> None:
    document = Document()
    document.add_heading(f"{report.metadata.organization} Intune Report", level=0)
    document.add_paragraph(f"Audience: {report.audience}")
    document.add_paragraph(f"Generated at: {report.metadata.generated_at}")
    document.add_page_break()
    _add_table_of_contents(document)

    assets_payload = next(
        (section.payload.get("assets") for section in report.sections if "assets" in section.payload),
        None,
    )

    for section in report.sections:
        document.add_page_break()
        document.add_heading(section.title, level=1)
        if section.description:
            document.add_paragraph(section.description)
        for key, payload in section.payload.items():
            if key == "summary":
                _render_summary_section(document, payload, assets_payload)
            elif key == "assets":
                _render_assets_section(document, payload)
            elif key == "assignment_coverage":
                _render_assignment_coverage_section(document, payload)
            else:
                payload_text = json.dumps(payload, indent=2, ensure_ascii=False)
                document.add_paragraph(payload_text)

    document.save(output_path)


def _add_table_of_contents(document: Document) -> None:
    document.add_heading("Table of Contents", level=1)
    paragraph = document.add_paragraph()
    run = paragraph.add_run()
    field = OxmlElement("w:fldSimple")
    field.set(qn("w:instr"), 'TOC \\o "1-3" \\h \\z \\u')
    run._r.append(field)


def _set_cell_shading(cell, color: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), color)
    tc_pr.append(shading)


def _render_summary_section(document: Document, payload: Dict[str, object], assets_payload: object) -> None:
    title = payload.get("title")
    if title:
        document.add_heading(str(title), level=2)
    highlights = payload.get("highlights", [])
    if highlights:
        document.add_paragraph("Highlights", style="List Bullet")
        for item in highlights:
            document.add_paragraph(str(item), style="List Bullet")

    metrics = payload.get("metrics", {})
    if metrics:
        document.add_paragraph("Key Metrics")
        table = document.add_table(rows=1, cols=2)
        table.style = "Light Grid"
        header_cells = table.rows[0].cells
        header_cells[0].text = "Metric"
        header_cells[1].text = "Value"
        _set_cell_shading(header_cells[0], "D9E1F2")
        _set_cell_shading(header_cells[1], "D9E1F2")
        for key, value in metrics.items():
            row_cells = table.add_row().cells
            row_cells[0].text = str(key)
            row_cells[1].text = json.dumps(value, ensure_ascii=False) if isinstance(value, (dict, list)) else str(value)

    if assets_payload:
        document.add_paragraph("Configuration Inventory")
        _render_configuration_inventory_table(document, assets_payload)
        document.add_paragraph("Platform Coverage")
        _render_platform_coverage_table(document, assets_payload)
        document.add_paragraph("Top Assigned Groups")
        _render_top_assigned_groups_table(document, assets_payload)


def _render_configuration_inventory_table(document: Document, assets_payload: object) -> None:
    inventory = _summarize_inventory(assets_payload)
    if not inventory:
        document.add_paragraph("No assets available.")
        return
    table = document.add_table(rows=1, cols=4)
    table.style = "Light Grid Accent 1"
    header_cells = table.rows[0].cells
    headers = ["Policy Type", "Total", "Assigned", "Unassigned"]
    for idx, header in enumerate(headers):
        header_cells[idx].text = header
        _set_cell_shading(header_cells[idx], "1F4E79")
    for row in inventory:
        row_cells = table.add_row().cells
        row_cells[0].text = row["asset_type"]
        row_cells[1].text = str(row["total"])
        row_cells[2].text = str(row["assigned"])
        row_cells[3].text = str(row["unassigned"])
        _set_cell_shading(row_cells[2], "C6E0B4")
        _set_cell_shading(row_cells[3], "FCE4D6")


def _render_platform_coverage_table(document: Document, assets_payload: object) -> None:
    coverage = _summarize_platform_coverage(assets_payload)
    if not coverage:
        document.add_paragraph("No platform data available.")
        return
    table = document.add_table(rows=1, cols=2)
    table.style = "Light Grid"
    header_cells = table.rows[0].cells
    header_cells[0].text = "Platform"
    header_cells[1].text = "Configs"
    _set_cell_shading(header_cells[0], "D9E1F2")
    _set_cell_shading(header_cells[1], "D9E1F2")
    for row in coverage:
        row_cells = table.add_row().cells
        row_cells[0].text = row["platform"]
        row_cells[1].text = str(row["count"])


def _render_top_assigned_groups_table(document: Document, assets_payload: object) -> None:
    group_summary = _summarize_groups(assets_payload)
    if not group_summary:
        document.add_paragraph("No group assignments available.")
        return
    top_groups = sorted(group_summary, key=lambda item: item["assigned_assets"], reverse=True)[:10]
    table = document.add_table(rows=1, cols=5)
    table.style = "Light Grid Accent 1"
    header_cells = table.rows[0].cells
    headers = ["Rank", "Group Name", "Assignments", "Settings Applied", "Type"]
    for idx, header in enumerate(headers):
        header_cells[idx].text = header
        _set_cell_shading(header_cells[idx], "BDD7EE")
    for idx, row in enumerate(top_groups, start=1):
        row_cells = table.add_row().cells
        row_cells[0].text = str(idx)
        row_cells[1].text = row["name"]
        row_cells[2].text = str(row["assigned_assets"])
        row_cells[3].text = str(row["settings_applied"])
        row_cells[4].text = row["type"]
        _set_cell_shading(row_cells[2], "C6E0B4")
        _set_cell_shading(row_cells[3], "FCE4D6")


def _summarize_groups(assets_payload: object) -> list[dict[str, object]]:
    summary: Dict[str, dict[str, object]] = {}
    if not isinstance(assets_payload, list):
        return []
    for asset in assets_payload:
        asset_id = asset.get("asset_id")
        assignments = asset.get("assignment_mappings", [])
        settings = asset.get("settings", {}) or {}
        settings_count = len(settings) if isinstance(settings, dict) else 0
        for mapping in assignments:
            group_id = mapping.get("groupId")
            if not group_id:
                continue
            group_entry = summary.setdefault(
                group_id,
                {
                    "name": mapping.get("groupDisplayName") or group_id,
                    "type": mapping.get("groupType") or "unknown",
                    "dynamic_rule": mapping.get("groupDynamicRule"),
                    "assigned_assets": set(),
                    "settings_applied": 0,
                },
            )
            group_entry["assigned_assets"].add(asset_id)
            group_entry["settings_applied"] += settings_count

    results: list[dict[str, object]] = []
    for group in summary.values():
        results.append(
            {
                "name": group["name"],
                "type": group["type"],
                "dynamic_rule": group["dynamic_rule"],
                "assigned_assets": len(group["assigned_assets"]),
                "settings_applied": group["settings_applied"],
            }
        )
    return sorted(results, key=lambda item: item["name"].lower())


def _summarize_inventory(assets_payload: object) -> list[dict[str, object]]:
    if not isinstance(assets_payload, list):
        return []
    summary: Dict[str, dict[str, int]] = {}
    for asset in assets_payload:
        asset_type = asset.get("asset_type") or "Unknown"
        assignments = asset.get("assignment_mappings", []) or asset.get("assignments", [])
        entry = summary.setdefault(str(asset_type), {"total": 0, "assigned": 0, "unassigned": 0})
        entry["total"] += 1
        if assignments:
            entry["assigned"] += 1
        else:
            entry["unassigned"] += 1
    results: list[dict[str, object]] = []
    for asset_type, counts in summary.items():
        results.append(
            {
                "asset_type": asset_type.replace("_", " ").title(),
                "total": counts["total"],
                "assigned": counts["assigned"],
                "unassigned": counts["unassigned"],
            }
        )
    return sorted(results, key=lambda item: item["asset_type"].lower())


def _extract_platforms(settings: Dict[str, object]) -> list[str]:
    platforms: list[str] = []
    if not isinstance(settings, dict):
        return platforms
    for key in ("platforms", "platform", "platformType"):
        value = settings.get(key)
        if isinstance(value, list):
            platforms.extend(str(item) for item in value if item)
        elif value:
            platforms.append(str(value))
    return platforms


def _summarize_platform_coverage(assets_payload: object) -> list[dict[str, object]]:
    if not isinstance(assets_payload, list):
        return []
    platform_counts: Dict[str, int] = {}
    for asset in assets_payload:
        settings = asset.get("settings", {}) or {}
        platforms = _extract_platforms(settings)
        for platform in platforms or ["Unknown"]:
            platform_counts[platform] = platform_counts.get(platform, 0) + 1
    return sorted(
        [{"platform": platform, "count": count} for platform, count in platform_counts.items()],
        key=lambda item: item["platform"].lower(),
    )


def _render_assets_section(document: Document, assets_payload: object) -> None:
    if not isinstance(assets_payload, list):
        document.add_paragraph("No asset data available.")
        return
    for asset in assets_payload:
        asset_name = asset.get("name") or "Unnamed Asset"
        asset_type = asset.get("asset_type") or "Unknown"
        document.add_heading(f"{asset_name} ({asset_type})", level=2)
        document.add_paragraph(f"Asset ID: {asset.get('asset_id')}")
        asset_description = asset.get("description")
        if asset_description:
            document.add_paragraph(str(asset_description))

        settings = asset.get("settings", {}) or {}
        document.add_paragraph("Settings")
        if settings:
            oma_rows = _extract_oma_setting_rows(settings)
            if oma_rows:
                _render_settings_table(document, oma_rows)
                remaining_settings = {
                    key: value for key, value in settings.items() if key != "settings"
                }
                if remaining_settings:
                    document.add_paragraph("Additional Settings")
                    _render_key_value_table(document, remaining_settings)
            else:
                _render_settings_table(document, _extract_setting_rows(settings))
        else:
            document.add_paragraph("No settings recorded.")

        assignments = asset.get("assignment_mappings", []) or []
        document.add_paragraph("Assignments")
        if assignments:
            table = document.add_table(rows=1, cols=7)
            table.style = "Light Grid"
            headers = [
                "Group",
                "Type",
                "Assignment",
                "Intent",
                "Delivery",
                "Schedule",
                "Dynamic Rule",
            ]
            header_cells = table.rows[0].cells
            for idx, header in enumerate(headers):
                header_cells[idx].text = header
                _set_cell_shading(header_cells[idx], "FFF2CC")
            for mapping in assignments:
                row_cells = table.add_row().cells
                row_cells[0].text = mapping.get("groupDisplayName") or mapping.get("groupId", "")
                row_cells[1].text = mapping.get("groupType") or ""
                row_cells[2].text = mapping.get("assignmentType") or ""
                row_cells[3].text = mapping.get("intent") or ""
                row_cells[4].text = mapping.get("delivery") or ""
                schedule = mapping.get("schedule")
                row_cells[5].text = (
                    json.dumps(schedule, ensure_ascii=False)
                    if isinstance(schedule, (dict, list))
                    else str(schedule or "")
                )
                row_cells[6].text = mapping.get("groupDynamicRule") or "N/A"
        else:
            document.add_paragraph("No assignments recorded.")


def _render_key_value_table(document: Document, settings: Dict[str, object]) -> None:
    table = document.add_table(rows=1, cols=2)
    table.style = "Light Grid"
    header_cells = table.rows[0].cells
    header_cells[0].text = "Setting"
    header_cells[1].text = "Value"
    _set_cell_shading(header_cells[0], "E2EFDA")
    _set_cell_shading(header_cells[1], "E2EFDA")
    for key, value in settings.items():
        row_cells = table.add_row().cells
        row_cells[0].text = str(key)
        row_cells[1].text = (
            json.dumps(value, ensure_ascii=False) if isinstance(value, (dict, list)) else str(value)
        )


def _render_settings_table(document: Document, rows: Iterable[dict[str, str]]) -> None:
    table = document.add_table(rows=1, cols=3)
    table.style = "Light Grid"
    header_cells = table.rows[0].cells
    header_cells[0].text = "Setting"
    header_cells[1].text = "Value"
    header_cells[2].text = "Description"
    _set_cell_shading(header_cells[0], "E2EFDA")
    _set_cell_shading(header_cells[1], "E2EFDA")
    _set_cell_shading(header_cells[2], "E2EFDA")
    for row in rows:
        row_cells = table.add_row().cells
        row_cells[0].text = row["setting"]
        row_cells[1].text = row["value"]
        row_cells[2].text = row["description"]


def _extract_oma_setting_rows(settings: Dict[str, object]) -> list[dict[str, str]]:
    raw_settings = settings.get("settings")
    if not isinstance(raw_settings, list):
        return []
    rows: list[dict[str, str]] = []
    for entry in raw_settings:
        if not isinstance(entry, dict):
            continue
        setting_name = (
            entry.get("omaUri")
            or entry.get("displayName")
            or entry.get("settingDefinitionId")
            or "Unnamed Setting"
        )
        description = (
            entry.get("description")
            or (entry.get("settingDefinition") or {}).get("description")
            or ""
        )
        value = entry.get("value")
        rows.append(
            {
                "setting": str(setting_name),
                "value": _stringify_setting_value(value),
                "description": str(description),
            }
        )
    return rows


def _stringify_setting_value(value: object) -> str:
    return json.dumps(value, ensure_ascii=False) if isinstance(value, (dict, list)) else str(value)


def _extract_setting_rows(settings: Dict[str, object]) -> list[dict[str, str]]:
    if not isinstance(settings, dict):
        return [{"setting": "N/A", "value": "N/A", "description": ""}]
    rows = _extract_oma_setting_rows(settings)
    if rows:
        remaining_settings = {key: value for key, value in settings.items() if key != "settings"}
    else:
        remaining_settings = settings
    for key, value in remaining_settings.items():
        rows.append({"setting": str(key), "value": _stringify_setting_value(value), "description": ""})
    return rows or [{"setting": "N/A", "value": "N/A", "description": ""}]


def _assignment_target_label(mapping: Dict[str, object]) -> str | None:
    target_type = str(mapping.get("targetType") or "").lower()
    if "alldevicesassignmenttarget" in target_type or "alldevices" in target_type:
        return "ALL DEVICES"
    if "alllicensedusersassignmenttarget" in target_type or "allusers" in target_type:
        return "ALL USERS"
    return None


def _render_assignment_coverage_section(document: Document, payload: Dict[str, object]) -> None:
    overview_table = document.add_table(rows=1, cols=2)
    overview_table.style = "Light Grid"
    header_cells = overview_table.rows[0].cells
    header_cells[0].text = "Metric"
    header_cells[1].text = "Value"
    _set_cell_shading(header_cells[0], "DDEBF7")
    _set_cell_shading(header_cells[1], "DDEBF7")
    for key in ("total_assets", "assigned_assets", "unassigned_assets"):
        row_cells = overview_table.add_row().cells
        row_cells[0].text = key.replace("_", " ").title()
        row_cells[1].text = str(payload.get(key, 0))

    assignments_by_group = payload.get("assignments_by_group", {}) or {}
    document.add_paragraph("Assignments by Group")
    if assignments_by_group:
        table = document.add_table(rows=1, cols=2)
        table.style = "Light Grid Accent 2"
        header_cells = table.rows[0].cells
        header_cells[0].text = "Group"
        header_cells[1].text = "Assigned Assets"
        _set_cell_shading(header_cells[0], "F8CBAD")
        _set_cell_shading(header_cells[1], "F8CBAD")
        for group, count in assignments_by_group.items():
            row_cells = table.add_row().cells
            row_cells[0].text = str(group)
            row_cells[1].text = str(count)
    else:
        document.add_paragraph("No group assignment coverage available.")


def _extract_assets_payload(report: RenderedReport) -> list[dict[str, object]]:
    assets_payload = next(
        (section.payload.get("assets") for section in report.sections if "assets" in section.payload),
        None,
    )
    if isinstance(assets_payload, list):
        return assets_payload
    return []


def _extract_assignment_coverage_payload(report: RenderedReport) -> Dict[str, object]:
    payload = next(
        (
            section.payload.get("assignment_coverage")
            for section in report.sections
            if "assignment_coverage" in section.payload
        ),
        None,
    )
    if isinstance(payload, dict):
        return payload
    return {}


def _extract_summary_payload(report: RenderedReport) -> Dict[str, object]:
    payload = next(
        (section.payload.get("summary") for section in report.sections if "summary" in section.payload),
        None,
    )
    if isinstance(payload, dict):
        return payload
    return {}


def _extract_setting_names(settings: Dict[str, object]) -> list[str]:
    names: list[str] = []
    if not isinstance(settings, dict):
        return ["N/A"]
    raw_settings = settings.get("settings")
    if isinstance(raw_settings, list):
        for entry in raw_settings:
            if not isinstance(entry, dict):
                continue
            setting_name = (
                entry.get("omaUri")
                or entry.get("displayName")
                or entry.get("settingDefinitionId")
                or "Unnamed Setting"
            )
            names.append(str(setting_name))
    for key in settings.keys():
        if key == "settings":
            continue
        names.append(str(key))
    return names or ["N/A"]


def _summarize_enrollment_profiles(
    assets_payload: list[dict[str, object]],
) -> list[dict[str, str]]:
    summary: Dict[str, list[str]] = {}
    for asset in assets_payload:
        if asset.get("asset_type") != "enrollment_profiles":
            continue
        profile_name = asset.get("name") or "Unnamed Enrollment Profile"
        for mapping in asset.get("assignment_mappings", []) or []:
            group_name = mapping.get("groupDisplayName") or mapping.get("groupId")
            if not group_name:
                continue
            summary.setdefault(str(group_name), []).append(str(profile_name))

    results: list[dict[str, str]] = []
    for group_name, profiles in summary.items():
        unique_profiles = sorted(set(profiles), key=str.lower)
        results.append(
            {
                "group": group_name,
                "profiles": ", ".join(unique_profiles),
            }
        )
    return sorted(results, key=lambda item: item["group"].lower())


def _write_pptx_report(report: RenderedReport, output_path: Path) -> None:
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = f"{report.metadata.organization} Intune Report"
    subtitle = title_slide.placeholders[1]
    subtitle.text = (
        f"Audience: {report.audience}\n"
        f"Generated: {report.metadata.generated_at}"
    )

    summary_payload = _extract_summary_payload(report)
    assets_payload = _extract_assets_payload(report)
    assignment_payload = _extract_assignment_coverage_payload(report)
    group_summary = _summarize_groups(assets_payload)
    enrollment_summary = _summarize_enrollment_profiles(assets_payload)
    assignments_by_group = assignment_payload.get("assignments_by_group", {}) if assignment_payload else {}

    summary_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    summary_slide.shapes.title.text = "Executive Summary"
    highlights = summary_payload.get("highlights", []) if isinstance(summary_payload, dict) else []
    metrics = summary_payload.get("metrics", {}) if isinstance(summary_payload, dict) else {}
    textbox = summary_slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(4.5), Inches(3.8))
    text_frame = textbox.text_frame
    text_frame.text = "Highlights"
    text_frame.paragraphs[0].font.bold = True
    for item in highlights:
        paragraph = text_frame.add_paragraph()
        paragraph.text = str(item)
        paragraph.level = 1

    metrics_table = summary_slide.shapes.add_table(
        rows=len(metrics) + 1,
        cols=2,
        left=Inches(5.2),
        top=Inches(1.4),
        width=Inches(4.0),
        height=Inches(3.0),
    ).table
    metrics_table.cell(0, 0).text = "Metric"
    metrics_table.cell(0, 1).text = "Value"
    for row_idx, (key, value) in enumerate(metrics.items(), start=1):
        metrics_table.cell(row_idx, 0).text = str(key).replace("_", " ").title()
        metrics_table.cell(row_idx, 1).text = str(value)

    policies_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    policies_slide.shapes.title.text = "Policies Applied to Groups"
    if group_summary:
        top_groups = sorted(group_summary, key=lambda item: item["assigned_assets"], reverse=True)[:10]
        chart_data = CategoryChartData()
        chart_data.categories = [row["name"] for row in top_groups]
        chart_data.add_series("Policies Applied", [row["assigned_assets"] for row in top_groups])
        chart = policies_slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.8),
            Inches(1.6),
            Inches(8.4),
            Inches(3.8),
            chart_data,
        ).chart
        chart.has_legend = False
        chart.value_axis.has_major_gridlines = True
        chart.plots[0].has_data_labels = True
    else:
        no_data_box = policies_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(1.0))
        no_data_box.text_frame.text = "No group assignment data available."

    relationship_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    relationship_slide.shapes.title.text = "Group Policy Relationships"
    if group_summary and assignments_by_group:
        top_groups = sorted(group_summary, key=lambda item: item["assigned_assets"], reverse=True)[:8]
        chart_data = BubbleChartData()
        series = chart_data.add_series("Group Impact")
        for row in top_groups:
            group_name = row["name"]
            series.add_data_point(
                row["assigned_assets"],
                row["settings_applied"],
                assignments_by_group.get(group_name, 0),
                group_name,
            )
        chart = relationship_slide.shapes.add_chart(
            XL_CHART_TYPE.BUBBLE,
            Inches(0.8),
            Inches(1.6),
            Inches(8.4),
            Inches(3.8),
            chart_data,
        ).chart
        chart.has_legend = False
        chart.value_axis.has_major_gridlines = True
        explanation_box = relationship_slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(8.4), Inches(0.4))
        explanation_box.text_frame.text = (
            "X = Policies Applied, Y = Settings Applied, Bubble Size = Group Assignment Count"
        )
    else:
        no_data_box = relationship_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(1.0))
        no_data_box.text_frame.text = "No group relationship data available."

    coverage_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    coverage_slide.shapes.title.text = "Group Coverage Overview"
    if assignments_by_group:
        sorted_groups = sorted(assignments_by_group.items(), key=lambda item: item[1], reverse=True)[:10]
        chart_data = CategoryChartData()
        chart_data.categories = [group for group, _ in sorted_groups]
        chart_data.add_series("Devices Applied", [count for _, count in sorted_groups])
        chart = coverage_slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.8),
            Inches(1.6),
            Inches(8.4),
            Inches(3.8),
            chart_data,
        ).chart
        chart.has_legend = False
        chart.plots[0].has_data_labels = True
    else:
        no_data_box = coverage_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(1.0))
        no_data_box.text_frame.text = "No coverage data available."

    enrollment_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    enrollment_slide.shapes.title.text = "Enrollment Profiles by Group"
    if enrollment_summary:
        table = enrollment_slide.shapes.add_table(
            rows=len(enrollment_summary) + 1,
            cols=2,
            left=Inches(0.6),
            top=Inches(1.6),
            width=Inches(9.0),
            height=Inches(4.0),
        ).table
        table.cell(0, 0).text = "Group"
        table.cell(0, 1).text = "Enrollment Profiles"
        for row_idx, row in enumerate(enrollment_summary, start=1):
            table.cell(row_idx, 0).text = row["group"]
            table.cell(row_idx, 1).text = row["profiles"]
    else:
        no_data_box = enrollment_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(1.0))
        no_data_box.text_frame.text = "No enrollment profile assignments available."

    presentation.save(output_path)


def _write_excel_report(report: RenderedReport, output_path: Path) -> None:
    workbook = Workbook()
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="4F81BD", end_color="4F81BD", fill_type="solid")
    section_title_font = Font(bold=True, size=14)
    border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )

    assets_payload = _extract_assets_payload(report)

    summary_sheet = workbook.active
    summary_sheet.title = "Summary"
    row_cursor = 1

    def add_section_title(sheet, title: str, row: int) -> int:
        cell = sheet.cell(row=row, column=1, value=title)
        cell.font = section_title_font
        return row + 1

    def add_table(sheet, headers: list[str], rows: list[list[object]], row: int) -> int:
        header_row = row
        for col_idx, header in enumerate(headers, start=1):
            cell = sheet.cell(row=header_row, column=col_idx, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = border
        row = header_row + 1
        for row_data in rows:
            for col_idx, value in enumerate(row_data, start=1):
                cell = sheet.cell(row=row, column=col_idx, value=value)
                cell.border = border
                cell.alignment = Alignment(horizontal="left")
            row += 1
        return row

    inventory_rows = [
        [row["asset_type"], row["total"], row["assigned"], row["unassigned"]]
        for row in _summarize_inventory(assets_payload)
    ]
    row_cursor = add_section_title(summary_sheet, "Configuration Inventory", row_cursor)
    row_cursor = add_table(
        summary_sheet,
        ["Policy Type", "Total", "Assigned", "Unassigned"],
        inventory_rows,
        row_cursor,
    )
    row_cursor += 2

    platform_rows = [
        [row["platform"], row["count"]]
        for row in _summarize_platform_coverage(assets_payload)
    ]
    row_cursor = add_section_title(summary_sheet, "Platform Coverage", row_cursor)
    row_cursor = add_table(summary_sheet, ["Platform", "Configs"], platform_rows, row_cursor)
    row_cursor += 2

    top_groups = sorted(_summarize_groups(assets_payload), key=lambda item: item["assigned_assets"], reverse=True)[:10]
    top_group_rows = [
        [idx, row["name"], row["assigned_assets"], row["settings_applied"], row["type"]]
        for idx, row in enumerate(top_groups, start=1)
    ]
    row_cursor = add_section_title(summary_sheet, "Top Assigned Groups", row_cursor)
    add_table(
        summary_sheet,
        ["Rank", "Group Name", "Assignments", "Settings Applied", "Type"],
        top_group_rows,
        row_cursor,
    )

    assignments_sheet = workbook.create_sheet("Assignments")
    assignments_headers = [
        "Setting",
        "Setting Value",
        "Description",
        "Policy",
        "Policy Description",
        "Policy Type",
        "Group",
        "Assigned Group Count",
        "Assignment Scope",
    ]
    for col_idx, header in enumerate(assignments_headers, start=1):
        cell = assignments_sheet.cell(row=1, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = border

    assignment_rows: list[list[object]] = []
    for asset in assets_payload:
        policy_name = asset.get("name") or "Unnamed Policy"
        policy_description = asset.get("description") or ""
        policy_type = asset.get("asset_type") or "Unknown"
        settings = asset.get("settings", {}) or {}
        assignments = asset.get("assignment_mappings", []) or []
        group_names = {
            mapping.get("groupDisplayName") or mapping.get("groupId")
            for mapping in assignments
            if mapping.get("groupDisplayName") or mapping.get("groupId")
        }
        group_names.discard(None)
        group_count = len(group_names)
        target_labels = {
            label for mapping in assignments if (label := _assignment_target_label(mapping))
        }
        if assignments:
            scope_parts = []
            if group_count:
                scope_parts.append("Groups")
            scope_parts.extend(sorted(target_labels))
            assignment_scope = ", ".join(scope_parts) if scope_parts else "Unassigned"
        else:
            assignment_scope = "Unassigned"

        group_labels = sorted(group_names, key=str.casefold) if group_names else []
        if not group_labels:
            group_labels = sorted(target_labels, key=str.casefold) if target_labels else ["Unassigned"]
        for setting_row in _extract_setting_rows(settings):
            for group_label in group_labels:
                assignment_rows.append(
                    [
                        setting_row["setting"],
                        setting_row["value"],
                        setting_row["description"],
                        policy_name,
                        policy_description,
                        policy_type,
                        group_label,
                        group_count,
                        assignment_scope,
                    ]
                )

    assignment_rows.sort(
        key=lambda row: (
            str(row[6]).casefold(),
            str(row[0]).casefold(),
            str(row[3]).casefold(),
        )
    )
    for row in assignment_rows:
        assignments_sheet.append(row)

    for row in assignments_sheet.iter_rows(
        min_row=2,
        max_row=assignments_sheet.max_row,
        max_col=assignments_sheet.max_column,
    ):
        for cell in row:
            cell.border = border
            if cell.column == 8:
                cell.alignment = Alignment(horizontal="center")
            else:
                cell.alignment = Alignment(horizontal="left")

    def autosize_columns(sheet) -> None:
        for col_idx in range(1, sheet.max_column + 1):
            column_letter = get_column_letter(col_idx)
            max_length = 0
            for cell in sheet[column_letter]:
                if cell.value is None:
                    continue
                max_length = max(max_length, len(str(cell.value)))
            sheet.column_dimensions[column_letter].width = min(max_length + 2, 50)

    autosize_columns(summary_sheet)
    autosize_columns(assignments_sheet)

    workbook.save(output_path)


def write_raw_export(raw_export: Dict[str, object], output_prefix: Path) -> Path:
    output_prefix.parent.mkdir(parents=True, exist_ok=True)
    output_path = output_prefix.with_name(f"{output_prefix.name}-raw.json")
    output_path.write_text(json.dumps(raw_export, indent=2, ensure_ascii=False), encoding="utf-8")
    return output_path
